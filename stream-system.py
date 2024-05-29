import streamlit as st
from pptx import Presentation
import requests
from io import BytesIO

# Define the function to download the PPTX file from a URL
def download_presentation(url):
    response = requests.get(url)
    response.raise_for_status()  # Check that the request was successful
    return BytesIO(response.content)

def load_presentation(file_path):
    return Presentation(file_path)

def display_presentation(presentation):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                st.write(shape.text)
            if shape.shape_type == 13:  # If shape is a picture
                image = shape.image
                with st.container():
                    st.image(image.blob)
            if shape.has_table:
                table = shape.table
                st.write("Table:")
                for row in table.rows:
                    st.write([cell.text for cell in row.cells])
            if shape.has_chart:
                chart = shape.chart
                # Streamlit cannot display the chart directly, but we can display data
                chart_data = chart.chart_data
                st.write("Chart data:")
                for series in chart_data.series:
                    st.write(f"Series {series.name}:")
                    for category, value in zip(chart_data.categories, series.values):
                        st.write(f"{category}: {value}")

# Apply custom CSS to change the background color and make font bold
st.markdown(
    """
    <style>
    .main {
        background-color: #656D4E;
        color: white;  /* Optional: Set text color to white for better contrast */
    }
    .main * {
        font-weight: bold;
    }
    </style>
    """,
    unsafe_allow_html=True
)

# Login credentials
USERNAME = "Ghaza"
PASSWORD = 123

# Create a login function
def login():
    st.title("Login")  # Display login title
    st.session_state['logged_in'] = False
    username = st.text_input("Username")
    password = st.text_input("Password", type="password")
    if st.button("Login"):
        if username == USERNAME and password == PASSWORD:
            st.session_state['logged_in'] = True
        else:
            st.error("Invalid username or password")

# Main function to run the app
def main():
    if 'logged_in' not in st.session_state:
        st.session_state['logged_in'] = False

    if not st.session_state['logged_in']:
        login()
    else:
        file_url = "https://github.com/MohammedHamza0/System/raw/main/System%20Overview%20(1)%20(1).pptx"  # Specify the file URL here
        file_path = download_presentation(file_url)
        presentation = load_presentation(file_path)
        display_presentation(presentation)

if __name__ == "__main__":
    main()







# import streamlit as st
# from pptx import Presentation
# import requests
# from io import BytesIO

# # Load presentation from URL
# def load_presentation(url):
#     try:
#         response = requests.get(url)
#         response.raise_for_status()  # Check for request errors
#         if response.headers.get('Content-Type') == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
#             pptx_file = BytesIO(response.content)
#             return Presentation(pptx_file)
#         else:
#             st.error("The downloaded file is not a valid PPTX file.")
#             st.error(f"Content-Type received: {response.headers.get('Content-Type')}")
#             # Save the invalid file for inspection
#             with open("invalid_file.pptx", "wb") as f:
#                 f.write(response.content)
#             st.error("The file has been saved as 'invalid_file.pptx' for further inspection.")
#     except requests.exceptions.RequestException as e:
#         st.error(f"Failed to download presentation: {e}")
#     except Exception as e:
#         st.error(f"Failed to load presentation: {e}")

# # Function to display the presentation
# def display_presentation(presentation):
#     for slide in presentation.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 st.write(shape.text)
#             if shape.shape_type == 13:  # Checking if shape is a picture
#                 image = shape.image
#                 with st.container():
#                     st.image(image.blob)

# # Apply custom CSS to change the background color and make font bold
# st.markdown(
#     """
#     <style>
#     .main {
#         background-color: #656D4E;
#         color: white;  /* Optional: Set text color to white for better contrast */
#     }
#     .main * {
#         font-weight: bold;
#     }
#     </style>
#     """,
#     unsafe_allow_html=True
# )

# # Credentials
# USERNAME = "Gaza"
# PASSWORD = "123"

# # Main function to run the app
# def main():
#     if 'logged_in' not in st.session_state:
#         st.session_state['logged_in'] = False

#     if not st.session_state['logged_in']:
#         login()
#     else:
#         pptx_url = "https://github.com/MohammedHamza0/System/blob/main/System%20Overview%20(1)%20(1).pptx"
#         presentation = load_presentation(pptx_url)
#         if presentation:
#             display_presentation(presentation)

# # Login function
# def login():
#     st.title("Login Page") 
#     username = st.text_input("Username")
#     password = st.text_input("Password", type="password")
#     if st.button("Login"):
#         if username == USERNAME and password == PASSWORD:
#             st.session_state['logged_in'] = True
#         else:
#             st.error("Invalid username or password")

# if __name__ == "__main__":
#     main()

